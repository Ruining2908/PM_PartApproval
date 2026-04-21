from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SCREENSHOT_DIR = ROOT / "artifacts" / "screenshots"
OUTPUT_PATH = ROOT / "artifacts" / "output" / "part-approval-user-guide.pptx"


SLIDES = [
    {
        "title": "Sign In To The Real Flutter App",
        "image": "real-01-login.png",
        "steps": [
            "Open the Part Approval web app built from the real Flutter code.",
            "Enter the password and click Login.",
            "After sign-in, the dashboard loads the request list automatically.",
        ],
    },
    {
        "title": "Review The Dashboard Overview",
        "image": "real-02-dashboard.png",
        "steps": [
            "Use the summary cards to monitor open requests and total approval value.",
            "Review the Requested By panel on the left to understand who submitted requests.",
            "Use the Part Request List on the right to review status chips for each request.",
        ],
    },
    {
        "title": "Search The Request List",
        "image": "real-03-search.png",
        "steps": [
            "Use the search bar to find requests by part name, machine, category, or requester.",
            "As you type, the list filters immediately without leaving the page.",
            "Use this search flow to narrow the queue before reviewing or approving a request.",
        ],
    },
    {
        "title": "Open Approval Detail",
        "image": "real-04-detail-modal.png",
        "steps": [
            "Select a request in the list to open the Approval Detail dialog.",
            "Review the request info, description, and remark before making a status change.",
            "Close the dialog to return to the list after your review is complete.",
        ],
    },
    {
        "title": "Confirm A Status Change",
        "image": "real-05-status-dialog.png",
        "steps": [
            "Use the status chips inside a request row to choose the next approval stage.",
            "When prompted, review the confirmation dialog carefully before proceeding.",
            "Click Confirm to apply the update, or Cancel to keep the current status unchanged.",
        ],
    },
]


def add_textbox(slide, left, top, width, height, text, size, bold=False, color=(31, 36, 48)):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_bullets(slide, left, top, width, height, lines):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, line in enumerate(lines):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = line
        p.level = 0
        p.bullet = True
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.size = Pt(18)
        run.font.color.rgb = RGBColor(61, 69, 88)
    return box


def add_scaled_image(slide, image_path, left, top, max_width, max_height):
    with Image.open(image_path) as img:
        width_px, height_px = img.size

    width_ratio = max_width / width_px
    height_ratio = max_height / height_px
    ratio = min(width_ratio, height_ratio)
    display_width = width_px * ratio
    display_height = height_px * ratio
    slide.shapes.add_picture(
        str(image_path),
        left,
        top,
        width=int(display_width),
        height=int(display_height),
    )


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_slide.background.fill.solid()
    title_slide.background.fill.fore_color.rgb = RGBColor(237, 239, 247)
    add_textbox(
        title_slide,
        Inches(0.7),
        Inches(0.9),
        Inches(8.8),
        Inches(0.8),
        "Part Approval Request User Guide",
        Pt(28),
        bold=True,
    )
    add_textbox(
        title_slide,
        Inches(0.7),
        Inches(1.85),
        Inches(4.8),
        Inches(2.2),
        "This presentation summarizes the real Flutter user workflow using browser-captured screenshots from the built application running in demo mode.",
        Pt(20),
    )
    add_scaled_image(
        title_slide,
        SCREENSHOT_DIR / "real-02-dashboard.png",
        Inches(6.7),
        Inches(0.7),
        Inches(5.8),
        Inches(5.8),
    )
    add_textbox(
        title_slide,
        Inches(0.7),
        Inches(5.9),
        Inches(5.5),
        Inches(0.6),
        "Screens included: login, dashboard overview, live search, approval detail, and status confirmation.",
        Pt(16),
        color=(96, 104, 124),
    )

    for slide_data in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(250, 251, 253)

        add_textbox(
            slide,
            Inches(0.55),
            Inches(0.35),
            Inches(5.4),
            Inches(0.5),
            slide_data["title"],
            Pt(24),
            bold=True,
        )
        add_textbox(
            slide,
            Inches(0.55),
            Inches(0.85),
            Inches(4.2),
            Inches(0.35),
            "Step-by-step reference for daily use",
            Pt(12),
            color=(96, 104, 124),
        )
        add_scaled_image(
            slide,
            SCREENSHOT_DIR / slide_data["image"],
            Inches(0.55),
            Inches(1.25),
            Inches(7.6),
            Inches(5.8),
        )
        add_bullets(
            slide,
            Inches(8.45),
            Inches(1.4),
            Inches(4.1),
            Inches(4.8),
            slide_data["steps"],
        )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved {OUTPUT_PATH}")


if __name__ == "__main__":
    build_deck()
